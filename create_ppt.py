from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create a presentation object
prs = Presentation()

# Define slide layouts
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

# Set a color scheme (Forest Green for brand)
brand_color = RGBColor(6, 95, 70)

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(title_slide_layout)
    title_box = slide.shapes.title
    subtitle_box = slide.placeholders[1]
    
    title_box.text = title
    subtitle_box.text = subtitle
    
    # Style title
    title_box.text_frame.paragraphs[0].font.color.rgb = brand_color
    title_box.text_frame.paragraphs[0].font.bold = True

def add_bullet_slide(title, points):
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.color.rgb = brand_color
    
    tf = body_shape.text_frame
    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20) # ensure text fits

# Slide 1: Title
add_title_slide(
    "Gokul Dairy Farm",
    "Digital Transformation Platform\nModern Cloud Solutions for Dairy Management"
)

# Slide 2: The Problem
add_bullet_slide(
    "The Problem: Manual Operations Limit Growth",
    [
        "Paper Ledgers: Tracking milk, deliveries, and skips on paper is error-prone.",
        "Billing Headaches: Calculating monthly bills manually takes hours.",
        "Lost Data: Physical notebooks can easily be lost or damaged.",
        "Lack of Insights: Difficult to track monthly revenue or herd performance."
    ]
)

# Slide 3: The Solution
add_bullet_slide(
    "The Solution: A Custom Cloud Platform",
    [
        "A lightning-fast, responsive web application for Mobile and Desktop.",
        "Centralized Data: Everything is stored securely in the cloud 24/7.",
        "Automated Math: Instantly calculates expected milk, skips, and total amounts.",
        "One-Click Actions: Send formatted bills and broadcast updates via WhatsApp."
    ]
)

# Slide 4: Key Features (Part 1)
add_bullet_slide(
    "Features: Dashboard & Daily Logging",
    [
        "Dashboard Analytics: See daily produced, delivered, and used milk at a glance.",
        "Revenue Tracking: 14-day line charts for milk and product sales.",
        "Daily Milk Log: Simple forms to log morning and evening yields.",
        "Skip Management: Easily pause a customer's delivery for a specific day."
    ]
)

# Slide 5: Key Features (Part 2)
add_bullet_slide(
    "Features: Smart WhatsApp Integration",
    [
        "1-Click Billing: Generates the exact bill (accounting for skipped days and custom rates).",
        "WhatsApp Sync: Tapping the WhatsApp button drafts a perfectly formatted message directly to the customer.",
        "Bulk Broadcast: Select multiple customers and instantly send a single WhatsApp announcement (e.g. price change) to all of them."
    ]
)

# Slide 6: Automated Billing & Payments
add_bullet_slide(
    "Features: Automated Billing & Payments",
    [
        "PDF Invoices: Automatically generates branded PDF bills for each customer.",
        "Payment Gateway: Integrated with Razorpay/Stripe.",
        "Instant Payments: Customers can pay directly via a secure link sent in WhatsApp.",
        "Automated Reconciliation: Invoices are automatically marked as 'Paid' in the database."
    ]
)

# Slide 7: Cross-Platform Availability
add_bullet_slide(
    "Cross-Platform Availability",
    [
        "Web Dashboard: Full administrative control from any PC or laptop.",
        "Mobile Apps: Fully published and available natively on the App Store (iOS) and Google Play (Android).",
        "Cloud-Synced: Data is perfectly synced across all devices instantly."
    ]
)

# Slide 8: Infrastructure & Security
add_bullet_slide(
    "Infrastructure & Enterprise Security",
    [
        "Cloud Hosting: Ultra-fast global hosting for zero downtime.",
        "Database (MongoDB Atlas): Secure cloud database capable of handling millions of records.",
        "Data Backups: Automated backups to ensure your data is never lost.",
        "Scalable API: Secure backend built to scale seamlessly as your farm grows."
    ]
)

# Save the presentation
output_path = r"C:\Users\Shubham\.gemini\antigravity\brain\f912fdf5-7f5d-4289-b816-e35074111bbe\scratch\Gokul_Dairy_Farm_Pitch.pptx"
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"Presentation saved successfully to {output_path}")
